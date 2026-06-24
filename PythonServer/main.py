from __future__ import annotations

import io
import json
import os
import re
from datetime import datetime
from typing import Any

import httpx
import qrcode
from fastapi import Body, FastAPI, Header, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pydantic import BaseModel, Field

# Load all configuration from a .env file next to this script, so settings live
# in one place instead of shell/$env exports. Must run before any os.environ.get.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ENV_PATH = os.path.join(_HERE, ".env")
try:
    from dotenv import load_dotenv

    load_dotenv(_ENV_PATH)
except ImportError:
    if os.path.exists(_ENV_PATH):
        print(
            "[warn] python-dotenv not installed; .env ignored. "
            "Run: pip install python-dotenv"
        )

API_BASE = os.environ.get(
    "BLUEJACK_API_BASE", "https://bluejack.binus.ac.id/lapi/api"
).rstrip("/")
SEMESTER_ID = os.environ.get(
    "SEMESTER_ID", "c93afa56-601f-46a6-ab51-671ba9a96328"
)
ACAD_BOOK_API = os.environ.get(
    "ACAD_BOOK_API",
    "https://academic-slc.apps.binus.ac.id/api/"
    "25b2ee66d775d76afebe917160f40c953760bc2dcfafefca5640daab9293a532/acad-book",
).rstrip("/")

_DEFAULT_PPT_DIR = os.path.normpath(os.path.join(_HERE, "..", "PPT"))
PPT_DIR = os.environ.get("PPT_TEMPLATE_DIR", _DEFAULT_PPT_DIR)

TEMPLATE_FILES = {
    "quiz": os.environ.get(
        "TEMPLATE_QUIZ",
        "02. [For Mhs] [Not For Share] - Briefing Pelaksanaan TM Even 2526.pptx",
    ),
    "uap": os.environ.get(
        "TEMPLATE_UAP",
        "02. [For Mhs] [Not For Share] - Briefing Pelaksanaan UAP dan Kumpul Proyek Even 2526.pptx",
    ),
}

HTTP_TIMEOUT = float(os.environ.get("HTTP_TIMEOUT", "30"))

SLIDE_HEADER_AND_QR_IDX = 1
SLIDE_TABLE_IDX = 2

# Defaults used in the assessment table when the course feed omits a field.
DEFAULT_DURATION = os.environ.get("BRIEFING_DEFAULT_DURATION", "100 menit di awal")
DEFAULT_PARTICIPANTS = os.environ.get("BRIEFING_DEFAULT_PARTICIPANTS", "1 orang")

# --- Testing overrides (real data only; no hardcoded course data) ----------
# Pretend "now" is a specific datetime so the server picks the REAL class that
# was scheduled at that moment. ISO format, e.g. BRIEFING_NOW_OVERRIDE=2026-03-12T11:30:00
NOW_OVERRIDE_RAW = os.environ.get("BRIEFING_NOW_OVERRIDE", "").strip()

# Which schedule window to query. Production wants "future"; to reach past
# sessions for testing set BRIEFING_SCHEDULE_MODE=all. When a NOW override is set
# we default to "all" so historical classes are reachable.
SCHEDULE_MODE = os.environ.get(
    "BRIEFING_SCHEDULE_MODE", "all" if NOW_OVERRIDE_RAW else "future"
).strip()

# Schedule source:
#   "classtransaction" - live ClassTransaction API, matched by assistant+day+shift
#   "api"              - live Schedule/GetJobsAssistant (dated assistant duties)
#   "file"             - TransactionSchedule.json (offline snapshot)
SCHEDULE_SOURCE = os.environ.get(
    "BRIEFING_SCHEDULE_SOURCE", "api"
).strip().lower()
TRANSACTION_FILE = os.environ.get(
    "BRIEFING_TRANSACTION_FILE",
    os.path.normpath(os.path.join(_HERE, "..", "TransactionSchedule.json")),
)
CLASS_TRANSACTION_URL = os.environ.get(
    "CLASS_TRANSACTION_URL",
    "https://bluejack.binus.ac.id/lapi/API/ClassTransaction/GetClassTransactionInSemester",
)
# Assistant code (e.g. "FB25-1") — normally comes from the card tap (request);
# this is only a fallback for manual testing without the desktop.
ASSISTANT_CODE = os.environ.get("BRIEFING_ASSISTANT_CODE", "").strip()
# Optional exact class filter (e.g. "BB07"). Normally the day+shift+assistant
# match is unique; the request's class_code (from the tapped room) takes priority.
CLASS_CODE = os.environ.get("BRIEFING_CLASS_CODE", "").strip()

# Day + shift to match. Hardcode for development (e.g. BRIEFING_DAY=4,
# BRIEFING_SHIFT=15:20-17:00); leave EMPTY for realtime (derived from the
# current weekday + time slot).
DAY_HARDCODE = os.environ.get("BRIEFING_DAY", "").strip()
SHIFT_HARDCODE = os.environ.get("BRIEFING_SHIFT", "").strip()

# Standard Binus practicum shift windows, used to derive the current shift
# in realtime mode.
SHIFTS = [
    "07:20-09:00",
    "09:20-11:00",
    "11:20-13:00",
    "13:20-15:00",
    "15:20-17:00",
    "17:20-19:00",
    "19:20-21:00",
]


def effective_now() -> datetime:
    """Return the override datetime when configured, else the real wall clock."""
    if NOW_OVERRIDE_RAW:
        try:
            return datetime.fromisoformat(NOW_OVERRIDE_RAW)
        except ValueError:
            pass
    return datetime.now()


def effective_day() -> int:
    """ISO weekday (Mon=1..Sun=7) to match against — hardcoded or realtime."""
    if DAY_HARDCODE:
        try:
            return int(DAY_HARDCODE)
        except ValueError:
            pass
    return effective_now().isoweekday()


def effective_shift() -> str | None:
    """Shift window (e.g. "15:20-17:00") to match — hardcoded, else the slot
    that contains the current time (None if outside all class hours)."""
    if SHIFT_HARDCODE:
        return SHIFT_HARDCODE
    now_t = effective_now().time()
    for shift in SHIFTS:
        m = _SHIFT_RE.search(shift)
        if not m:
            continue
        sh, sm, eh, em = (int(g) for g in m.groups())
        if datetime.min.replace(hour=sh, minute=sm).time() <= now_t <= datetime.min.replace(
            hour=eh, minute=em
        ).time():
            return shift
    return None



class GenerateRequest(BaseModel):
    line_group_link: str = Field(..., description="URL of the LINE group to encode in the QR")
    template_type: str = Field(..., description='"quiz" or "uap"')
    # The assistant initial (e.g. "FB25-1") comes from the card tap via NuxtJS.
    # Falls back to BRIEFING_ASSISTANT_CODE in .env when omitted (testing only).
    assistant_code: str | None = Field(None, description="Assistant initial from the card tap")
    # Optional class (e.g. "BB07") to disambiguate same-slot classes; in future
    # this comes from the tapped room. Falls back to BRIEFING_CLASS_CODE.
    class_code: str | None = Field(None, description="Class code to disambiguate")



async def fetch_schedule(token: str) -> list[dict]:
    url = f"{API_BASE}/Schedule/GetJobsAssistant"
    params = {"mode": SCHEDULE_MODE, "semesterId": SEMESTER_ID}
    headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
    async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as client:
        r = await client.get(url, params=params, headers=headers)
    if r.status_code == 401:
        raise HTTPException(401, "Bluejack rejected the bearer token")
    if r.status_code >= 400:
        raise HTTPException(
            502, f"Bluejack Schedule API returned {r.status_code}: {r.text[:200]}"
        )
    data = r.json()
    if isinstance(data, dict):
        for v in data.values():
            if isinstance(v, list):
                return v
        raise HTTPException(502, "Schedule API did not return a list")
    return data


async def fetch_class_transactions(token: str) -> list[dict]:
    headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
    async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as client:
        r = await client.get(
            CLASS_TRANSACTION_URL,
            params={"semesterId": SEMESTER_ID},
            headers=headers,
        )
    if r.status_code == 401:
        raise HTTPException(401, "Bluejack rejected the bearer token")
    if r.status_code >= 400:
        raise HTTPException(
            502, f"ClassTransaction API returned {r.status_code}: {r.text[:200]}"
        )
    data = r.json()
    if isinstance(data, list):
        return data
    if isinstance(data, dict):
        for v in data.values():
            if isinstance(v, list):
                return v
    raise HTTPException(502, "ClassTransaction API did not return a list")


async def fetch_course_data(course_code: str) -> dict:
    url = f"{ACAD_BOOK_API}/{course_code}"
    async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as client:
        r = await client.post(url)
    if r.status_code >= 400:
        raise HTTPException(
            502, f"acad-book API returned {r.status_code}: {r.text[:200]}"
        )
    return r.json()


def find_current_job(jobs: list[dict]) -> dict | None:
    now = effective_now()
    for job in jobs:
        start_raw = job.get("StartDate")
        end_raw = job.get("EndDate")
        if not start_raw or not end_raw:
            continue
        try:
            start = datetime.fromisoformat(start_raw)
            end = datetime.fromisoformat(end_raw)
        except ValueError:
            continue
        if start <= now <= end:
            return job
    return None


def load_transaction_schedule() -> list[dict]:
    with open(TRANSACTION_FILE, encoding="utf-8") as f:
        data = json.load(f)
    return data if isinstance(data, list) else []


# "Day" in TransactionSchedule is ISO weekday (Mon=1 .. Sun=7), confirmed against
# real dates (BB07 is Day 1, and 30-Mar-2026 — a Monday — was a BB07 session).
_SHIFT_RE = re.compile(r"(\d{1,2}):(\d{2})\s*-\s*(\d{1,2}):(\d{2})")


def find_current_transaction(
    entries: list[dict], assistant_code: str, class_code: str = ""
) -> dict | None:
    now = effective_now()
    weekday = now.isoweekday()
    for e in entries:
        if e.get("Day") != weekday:
            continue
        m = _SHIFT_RE.search(e.get("Shift") or "")
        if not m:
            continue
        sh, sm, eh, em = (int(g) for g in m.groups())
        start = now.replace(hour=sh, minute=sm, second=0, microsecond=0)
        end = now.replace(hour=eh, minute=em, second=0, microsecond=0)
        if not (start <= now <= end):
            continue
        if assistant_code and assistant_code.lower() not in (
            e.get("Assistant") or ""
        ).lower():
            continue
        if class_code and class_code.lower() != (e.get("Class") or "").strip().lower():
            continue
        return e
    return None


def find_class_transaction(
    entries: list[dict],
    assistant_code: str,
    day: int,
    shift: str | None,
    class_code: str = "",
) -> dict | None:
    """Filter the live ClassTransaction list by assistant + day + shift (and
    optionally class). The class code is read off the matched entry."""
    for e in entries:
        if e.get("Day") != day:
            continue
        if shift and (e.get("Shift") or "").strip() != shift.strip():
            continue
        if assistant_code and assistant_code.lower() not in (
            e.get("Assistant") or ""
        ).lower():
            continue
        if class_code and class_code.lower() != (e.get("Class") or "").strip().lower():
            continue
        return e
    return None


_CLASS_CODE_RE = re.compile(r"^[A-Z]{2,}\d{2,}$")


def parse_description(description: str) -> tuple[str, str, str]:
    if "-" not in description:
        raise HTTPException(500, f"Unparseable description: {description!r}")
    kd_mtk, rest = description.split("-", 1)
    kd_mtk = kd_mtk.strip()
    tokens = rest.strip().split()
    if not tokens:
        raise HTTPException(500, f"Description has no course name: {description!r}")

    class_idx: int | None = None
    for i in range(len(tokens) - 1, -1, -1):
        if _CLASS_CODE_RE.match(tokens[i]):
            class_idx = i
            break
    if class_idx is None:
        class_idx = len(tokens) - 1

    kd_kelas = tokens[class_idx].strip()
    nama_mtk = " ".join(tokens[:class_idx]).strip()
    return kd_mtk, nama_mtk, kd_kelas


def parse_notes(notes: list[str]) -> dict[str, str]:
    out = {"duration": "", "participants": "", "max_file_size": ""}
    for note in notes or []:
        low = note.lower()

        m = re.search(r"held for\s+(\d+)\s*minutes?", low)
        if m and not out["duration"]:
            out["duration"] = f"{m.group(1)} menit di awal"

        m = re.search(r"each group contains\s+(\d+)\s*-\s*(\d+)\s+members?", low)
        if m and not out["participants"]:
            out["participants"] = f"{m.group(1)} - {m.group(2)} orang"
        else:
            m = re.search(r"each group contains\s+(\d+)\s+members?", low)
            if m and not out["participants"]:
                out["participants"] = f"{m.group(1)} orang"

        m = re.search(r"maximum file size\s+(\d+)\s*(mb|gb|kb)\b", low)
        if m and not out["max_file_size"]:
            out["max_file_size"] = f"{m.group(1)} {m.group(2).upper()}"
    return out


def pick_info_entry(info_list: list[dict], template_type: str) -> dict | None:
    if not info_list:
        return None
    if template_type == "quiz":
        return min(info_list, key=lambda x: x.get("meeting", 10**9))
    return max(info_list, key=lambda x: x.get("meeting", -1))



def generate_qr_png(text: str) -> bytes:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=10,
        border=2,
    )
    qr.add_data(text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


_TOKEN_RUNS = {"KdMtk", "NamaMtk", "KdKelas"}


def replace_placeholders_in_text_frame(
    text_frame, run_map: dict[str, str], substring_map: dict[str, str]
) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            if text in run_map:
                run.text = run_map[text]
                continue
            for needle, value in substring_map.items():
                if needle in text:
                    run.text = text.replace(needle, value)


def replace_qr_picture(slide, qr_png_bytes: bytes) -> bool:
    target = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            target = shape
            break
    if target is None:
        return False

    left, top, width, height = target.left, target.top, target.width, target.height

    sp = target._element
    sp.getparent().remove(sp)

    slide.shapes.add_picture(
        io.BytesIO(qr_png_bytes), left, top, width=width, height=height
    )
    return True


def set_cell_text_keep_format(cell, new_text: str) -> None:
    tf = cell.text_frame
    paragraphs = tf.paragraphs
    if not paragraphs:
        cell.text = new_text
        return

    first_p = paragraphs[0]
    p_elem = first_p._p

    txBody = p_elem.getparent()
    for extra in list(txBody.findall(qn("a:p")))[1:]:
        txBody.remove(extra)

    runs = first_p.runs
    if not runs:
        first_p.text = new_text
        return

    for extra_r in list(p_elem.findall(qn("a:r")))[1:]:
        p_elem.remove(extra_r)

    for tag in ("a:br", "a:fld"):
        for el in list(p_elem.findall(qn(tag))):
            p_elem.remove(el)

    runs[0].text = new_text


def _find_main_table(slide):
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        first_col = [r.cells[0].text.strip() for r in shape.table.rows]
        if any(label in first_col for label in ("TM/Quiz", "UAP", "Proyek")):
            return shape.table
    return None


def _fill_component_row(
    table, find_label: str, info: dict, new_label: str | None = None
) -> None:
    target_row = None
    for row in table.rows:
        if find_label in row.cells[0].text:
            target_row = row
            break
    if target_row is None:
        return

    parsed = parse_notes(info.get("notes", []))
    meeting = info.get("meeting")
    cells = list(target_row.cells)

    # Optionally rename the row label (e.g. "TM/Quiz" -> "TM/Quiz 1").
    if new_label:
        set_cell_text_keep_format(cells[0], new_label)

    set_cell_text_keep_format(cells[2], str(meeting) if meeting is not None else "-")
    # Participants/duration fall back to a default when the feed omits them
    # (max file size keeps a dash — there's no sensible default).
    set_cell_text_keep_format(cells[3], parsed["participants"] or DEFAULT_PARTICIPANTS)
    set_cell_text_keep_format(cells[4], parsed["duration"] or DEFAULT_DURATION)
    set_cell_text_keep_format(cells[6], parsed["max_file_size"] or "-")


def _is_exam(info: dict | None) -> bool:
    comp = ((info or {}).get("component") or "").lower()
    return "final" in comp or "exam" in comp or "uap" in comp


def populate_main_table(slide, quiz_info: dict | None, uap_info: dict | None) -> None:
    """Fill both assessment rows so the slide shows the full grading breakdown.
    The Proyek row is left untouched (no project data in the course feed).

    When a course has two non-exam components (e.g. Assignment 1 + Assignment 2,
    no Final Exam) it's two quizzes, not quiz + UAP — so the rows are relabelled
    "TM/Quiz 1" and "TM/Quiz 2"."""
    table = _find_main_table(slide)
    if table is None:
        return

    two_quizzes = (
        quiz_info is not None
        and uap_info is not None
        and quiz_info is not uap_info
        and not _is_exam(quiz_info)
        and not _is_exam(uap_info)
    )

    if quiz_info:
        _fill_component_row(
            table, "TM/Quiz", quiz_info, "TM/Quiz 1" if two_quizzes else None
        )
    if uap_info:
        _fill_component_row(
            table, "UAP", uap_info, "TM/Quiz 2" if two_quizzes else None
        )


def populate_software_extension_table(
    slide, software: list[dict], extension: str
) -> None:
    target = None
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        first_col = [r.cells[0].text.strip() for r in shape.table.rows]
        if "Software" in first_col and "Extension" in first_col:
            target = shape.table
            break
    if target is None:
        return

    software_str = ", ".join(
        str(s.get("name", "")).strip()
        for s in (software or [])
        if s.get("name")
    )

    for row in target.rows:
        label = row.cells[0].text.strip()
        if label == "Software" and software_str:
            set_cell_text_keep_format(row.cells[1], software_str)
        elif label == "Extension" and extension:
            set_cell_text_keep_format(row.cells[1], extension)


def build_pptx(
    template_path: str,
    template_type: str,
    kd_mtk: str,
    nama_mtk: str,
    kd_kelas: str,
    line_group_link: str,
    course: dict[str, Any],
) -> bytes:
    prs = Presentation(template_path)

    run_map = {"KdMtk": kd_mtk, "NamaMtk": nama_mtk, "KdKelas": kd_kelas}
    substring_map = {"[Link Group]": line_group_link}

    if len(prs.slides) <= SLIDE_TABLE_IDX:
        raise HTTPException(
            500, f"Template only has {len(prs.slides)} slides; expected >= 3"
        )

    slide2 = prs.slides[SLIDE_HEADER_AND_QR_IDX]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            replace_placeholders_in_text_frame(
                shape.text_frame, run_map, substring_map
            )

    qr_png = generate_qr_png(line_group_link)
    replace_qr_picture(slide2, qr_png)

    slide3 = prs.slides[SLIDE_TABLE_IDX]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            replace_placeholders_in_text_frame(
                shape.text_frame, run_map, substring_map
            )

    info_list = course.get("info", []) or []
    software_list = course.get("software", []) or []
    quiz_info = pick_info_entry(info_list, "quiz")
    uap_info = pick_info_entry(info_list, "uap")

    populate_main_table(slide3, quiz_info, uap_info)

    # Software/extension follow the requested template's primary assessment.
    primary_info = quiz_info if template_type == "quiz" else uap_info
    extension = (primary_info or {}).get("extension", "") or ""
    populate_software_extension_table(slide3, software_list, extension)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

app = FastAPI(title="Briefing PPTX Generator")


@app.get("/health")
async def health() -> dict[str, Any]:
    templates_ok = {
        k: os.path.exists(os.path.join(PPT_DIR, v))
        for k, v in TEMPLATE_FILES.items()
    }
    return {
        "status": "ok",
        "ppt_dir": PPT_DIR,
        "templates": templates_ok,
        "schedule_source": SCHEDULE_SOURCE,
        "schedule_mode": SCHEDULE_MODE,
        "assistant_code": ASSISTANT_CODE or None,
        "class_code": CLASS_CODE or None,
        "day_hardcode": DAY_HARDCODE or None,
        "shift_hardcode": SHIFT_HARDCODE or None,
        "effective_day": effective_day(),
        "effective_shift": effective_shift(),
        "now_override": NOW_OVERRIDE_RAW or None,
        "effective_now": effective_now().isoformat(),
        "now": datetime.now().isoformat(),
    }


@app.post("/generate-briefing")
async def generate_briefing(
    req: GenerateRequest = Body(...),
    authorization: str = Header(..., description="Bearer <bluejack-token>"),
):
    if not authorization.lower().startswith("bearer "):
        raise HTTPException(401, "Authorization header must be 'Bearer <token>'")
    token = authorization.split(" ", 1)[1].strip()
    if not token:
        raise HTTPException(401, "Empty bearer token")

    template_type = req.template_type.strip().lower()
    if template_type not in TEMPLATE_FILES:
        raise HTTPException(
            400, f"template_type must be 'quiz' or 'uap', got {req.template_type!r}"
        )

    template_path = os.path.join(PPT_DIR, TEMPLATE_FILES[template_type])
    if not os.path.exists(template_path):
        raise HTTPException(
            500, f"Template file missing on server: {template_path}"
        )

    # Assistant initial + class come from the card tap (via the request); the
    # .env values are only a fallback for local testing.
    assistant = (req.assistant_code or ASSISTANT_CODE or "").strip()
    class_code = (req.class_code or CLASS_CODE or "").strip()

    if SCHEDULE_SOURCE == "classtransaction":
        entries = await fetch_class_transactions(token)
        day = effective_day()
        shift = effective_shift()
        current = find_class_transaction(entries, assistant, day, shift, class_code)
        if not current:
            raise HTTPException(
                404,
                "No class transaction matched "
                f"(assistant={assistant or 'any'}, day={day}, shift={shift or 'any'}, "
                f"class={class_code or 'any'})",
            )
        description = (
            f"{current.get('Subject', '')} {(current.get('Class') or '').strip()}"
        )
    elif SCHEDULE_SOURCE == "file":
        entries = load_transaction_schedule()
        current = find_current_transaction(entries, assistant, class_code)
        if not current:
            raise HTTPException(
                404,
                "No class in TransactionSchedule for "
                f"{effective_now().isoformat()} "
                f"(assistant={assistant or 'any'}, class={class_code or 'any'})",
            )
        # Build a description from Subject + Class, e.g.
        # "COMP6799001-Database Technology" + "BB07" -> reuse parse_description.
        description = (
            f"{current.get('Subject', '')} {(current.get('Class') or '').strip()}"
        )
    else:
        jobs = await fetch_schedule(token)
        current = find_current_job(jobs)
        if not current:
            raise HTTPException(
                404,
                "No scheduled class is active "
                f"(checked {len(jobs)} jobs for {effective_now().isoformat()})",
            )
        description = current.get("Description", "")

    kd_mtk, nama_mtk, kd_kelas = parse_description(description)

    payload = await fetch_course_data(kd_mtk)
    try:
        course = payload["payloads"]["acad-book"]["course"][0]
    except (KeyError, IndexError, TypeError):
        raise HTTPException(
            502,
            f"acad-book response missing payloads.acad-book.course[0] for {kd_mtk}",
        )

    pptx_bytes = build_pptx(
        template_path=template_path,
        template_type=template_type,
        kd_mtk=kd_mtk,
        nama_mtk=nama_mtk,
        kd_kelas=kd_kelas,
        line_group_link=req.line_group_link,
        course=course,
    )

    filename = f"Briefing_{kd_mtk}_{kd_kelas}_{template_type}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "X-Course-Code": kd_mtk,
            "X-Class-Code": kd_kelas,
        },
    )


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(
        "main:app",
        host=os.environ.get("HOST", "0.0.0.0"),
        port=int(os.environ.get("PORT", "8000")),
        reload=bool(os.environ.get("RELOAD")),
    )
